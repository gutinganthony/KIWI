"""Generate 5-slide PPTX: KIWI × Claude dashboard intro."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── Palette (TradingView dark) ──
BG_DARK   = RGBColor(0x13, 0x17, 0x22)   # #131722
BG_CARD   = RGBColor(0x1E, 0x22, 0x2D)   # #1e222d
BG_LIGHT  = RGBColor(0x2A, 0x2E, 0x39)   # #2a2e39
WHITE     = RGBColor(0xD1, 0xD4, 0xDC)   # text
MUTED     = RGBColor(0x78, 0x7B, 0x86)
GREEN     = RGBColor(0x26, 0xA6, 0x9A)
YELLOW    = RGBColor(0xF7, 0xC9, 0x48)
ORANGE    = RGBColor(0xFF, 0x98, 0x00)
RED       = RGBColor(0xEF, 0x53, 0x50)
BLUE      = RGBColor(0x42, 0xA5, 0xF5)
PURPLE    = RGBColor(0x7C, 0x4D, 0xFF)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helper: fill slide background ──
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Helper: add text box ──
def tb(slide, text, x, y, w, h,
       size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
       wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

# ── Helper: filled rectangle ──
def rect(slide, x, y, w, h, color, radius=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 workaround
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# Use the proper enum
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

def box(slide, x, y, w, h, fill_color, line_color=None):
    from pptx.oxml.ns import qn
    sp = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    return sp

def accent_bar(slide, x, y, h, color):
    box(slide, x, y, Inches(0.08), h, color)

# ── SLIDE 1: Title ──────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
bg(s1, BG_DARK)

# gradient-ish accent band on left
box(s1, Inches(0), Inches(0), Inches(0.5), H, BG_CARD)
box(s1, Inches(0), Inches(0), Inches(0.12), H, BLUE)

# top-right tag
tb(s1, "2026 · Personal Project", Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.5),
   size=11, color=MUTED, align=PP_ALIGN.RIGHT)

# title
tb(s1, "KIWI × Claude", Inches(0.9), Inches(1.6), Inches(11), Inches(1.2),
   size=52, bold=True, color=WHITE)

# subtitle
tb(s1, "三合一市場預警 Dashboard", Inches(0.9), Inches(2.85), Inches(11), Inches(0.8),
   size=32, bold=False, color=BLUE)

# divider
box(s1, Inches(0.9), Inches(3.75), Inches(4.5), Inches(0.05), BLUE)

# three badge boxes
badges = [
    ("AVI V5", "估值環境", YELLOW),
    ("CPI",    "崩盤概率", GREEN),
    ("TSI",    "科技壓力", ORANGE),
]
for i, (tag, desc, col) in enumerate(badges):
    bx = Inches(0.9) + i * Inches(2.0)
    box(s1, bx, Inches(4.1), Inches(1.7), Inches(0.7), BG_CARD)
    tb(s1, tag,  bx + Inches(0.1), Inches(4.15), Inches(1.5), Inches(0.45),
       size=20, bold=True, color=col)
    tb(s1, desc, bx + Inches(0.1), Inches(4.55), Inches(1.5), Inches(0.3),
       size=11, color=MUTED)

# bottom tagline
tb(s1, "從原始碼到 GitHub Pages — 一個 Claude Session 完成",
   Inches(0.9), Inches(5.7), Inches(10), Inches(0.5),
   size=15, italic=True, color=MUTED)

# GitHub Pages badge
tb(s1, "🌐  gutinganthony.github.io/KIWI",
   Inches(0.9), Inches(6.3), Inches(10), Inches(0.5),
   size=13, color=BLUE)


# ── SLIDE 2: The Problem ────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
bg(s2, BG_DARK)
box(s2, Inches(0), Inches(0), Inches(0.12), H, YELLOW)

tb(s2, "我們在解決什麼問題？", Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
   size=30, bold=True, color=WHITE)
tb(s2, "市場風險很難用一個數字看懂", Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
   size=16, color=MUTED)
box(s2, Inches(0.5), Inches(1.6), Inches(4.5), Inches(0.05), YELLOW)

# Left column: OLD WAY
box(s2, Inches(0.5), Inches(1.9), Inches(3.8), Inches(4.8), BG_CARD)
accent_bar(s2, Inches(0.5), Inches(1.9), Inches(4.8), RED)
tb(s2, "❌  以前的做法", Inches(0.7), Inches(2.0), Inches(3.5), Inches(0.5),
   size=15, bold=True, color=RED)
old_items = [
    "估值指標（CAPE）每月手算",
    "VIX / 殖利率各看各的",
    "崩盤信號只能事後知道",
    "科技股壓力無法量化",
    "三個工具散在不同地方",
    "手機完全沒辦法看",
]
for j, it in enumerate(old_items):
    tb(s2, f"  ·  {it}", Inches(0.7), Inches(2.55) + j * Inches(0.55),
       Inches(3.4), Inches(0.5), size=13, color=MUTED)

# Right column: NEW WAY
box(s2, Inches(5.2), Inches(1.9), Inches(7.6), Inches(4.8), BG_CARD)
accent_bar(s2, Inches(5.2), Inches(1.9), Inches(4.8), GREEN)
tb(s2, "✅  現在的做法", Inches(5.4), Inches(2.0), Inches(7), Inches(0.5),
   size=15, bold=True, color=GREEN)
new_items = [
    ("AVI V5",  "月度估值環境評分（0-10）",          YELLOW),
    ("CPI",     "日頻崩盤概率（0-100）23-42天預警",   GREEN),
    ("TSI",     "日頻科技股壓力（0-100）",            ORANGE),
    ("一頁搞定", "三系統整合、綜合信號、行動建議",     BLUE),
    ("手機可用", "GitHub Pages，隨時查看",            WHITE),
    ("雙語介面", "EN / 中文一鍵切換",                 MUTED),
]
for j, (tag, desc, col) in enumerate(new_items):
    ty = Inches(5.4)
    y  = Inches(2.55) + j * Inches(0.55)
    tb(s2, tag,  ty, y, Inches(1.4), Inches(0.5), size=13, bold=True, color=col)
    tb(s2, desc, ty + Inches(1.45), y, Inches(5.5), Inches(0.5), size=13, color=MUTED)


# ── SLIDE 3: Three Indexes ──────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
bg(s3, BG_DARK)
box(s3, Inches(0), Inches(0), Inches(0.12), H, GREEN)

tb(s3, "三個指數各自做什麼？", Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
   size=30, bold=True, color=WHITE)
tb(s3, "互相補充，不互相取代", Inches(0.5), Inches(1.05), Inches(12), Inches(0.5),
   size=15, italic=True, color=MUTED)
box(s3, Inches(0.5), Inches(1.55), Inches(4), Inches(0.05), GREEN)

cards = [
    {
        "tag": "AVI V5",
        "full": "Adjusted Valuation Index",
        "zh": "估值環境指數",
        "color": YELLOW,
        "freq": "月度更新",
        "range": "0 – 10 分",
        "question": "「市場現在有多貴？」",
        "items": [
            "14 個指標 × 6 個維度",
            "20 年滾動百分位標準化",
            "估值 / 信用 / 宏觀 / 利率",
            "動量 / 獲利能力",
            "5/5 歷史危機偵測成功",
            "當前：4.5/10 🟡 偏高",
        ],
        "stat": "100%",
        "stat_label": "5次危機偵測率",
    },
    {
        "tag": "CPI",
        "full": "Crash Probability Index",
        "zh": "崩盤概率指數",
        "color": GREEN,
        "freq": "日頻更新",
        "range": "0 – 100 分",
        "question": "「接下來一個月會崩嗎？」",
        "items": [
            "12 個日頻指標",
            "VIX 結構 / 信用加速 / 廣度",
            "RSI 背離 / 動量崩潰",
            "殖利率急升 / 盤中賣壓",
            "Flash Alert（4+ 指標同觸）",
            "當前：23/100 🟢 正常",
        ],
        "stat": "28天",
        "stat_label": "平均預警提前期",
    },
    {
        "tag": "TSI",
        "full": "Tech Stress Index",
        "zh": "科技壓力指數",
        "color": ORANGE,
        "freq": "日頻更新",
        "range": "0 – 100 分",
        "question": "「科技股現在承受多大壓力？」",
        "items": [
            "9 個市場代理指標",
            "SOX/QQQ 背離",
            "殖利率衝擊（10Y / 30Y）",
            "熊市陡化 / 記憶體動能",
            "VIX-科技相關性",
            "當前：45/100 🟡 謹慎",
        ],
        "stat": "7/7",
        "stat_label": "壓力事件偵測率",
    },
]

for i, c in enumerate(cards):
    cx = Inches(0.45) + i * Inches(4.28)
    cw = Inches(4.1)
    ch = Inches(5.5)
    cy = Inches(1.75)
    box(s3, cx, cy, cw, ch, BG_CARD)
    accent_bar(s3, cx, cy, ch, c["color"])

    tb(s3, c["tag"],  cx + Inches(0.2), cy + Inches(0.12), cw - Inches(0.3), Inches(0.5),
       size=22, bold=True, color=c["color"])
    tb(s3, c["zh"],   cx + Inches(0.2), cy + Inches(0.58), cw - Inches(0.3), Inches(0.35),
       size=13, color=WHITE)
    tb(s3, c["full"], cx + Inches(0.2), cy + Inches(0.9),  cw - Inches(0.3), Inches(0.3),
       size=10, italic=True, color=MUTED)

    # freq + range chips
    tb(s3, c["freq"],  cx + Inches(0.2), cy + Inches(1.25), Inches(1.5), Inches(0.3),
       size=10, color=c["color"], bold=True)
    tb(s3, c["range"], cx + Inches(1.8), cy + Inches(1.25), Inches(1.5), Inches(0.3),
       size=10, color=MUTED)

    box(s3, cx + Inches(0.18), cy + Inches(1.6), cw - Inches(0.36), Inches(0.03), BG_LIGHT)

    tb(s3, c["question"], cx + Inches(0.2), cy + Inches(1.7), cw - Inches(0.3), Inches(0.45),
       size=11, italic=True, color=WHITE)

    for k, it in enumerate(c["items"]):
        tb(s3, f"  ·  {it}", cx + Inches(0.18), cy + Inches(2.2) + k * Inches(0.42),
           cw - Inches(0.3), Inches(0.4), size=11, color=MUTED)

    # stat badge
    box(s3, cx + Inches(0.18), cy + ch - Inches(0.72), cw - Inches(0.36), Inches(0.6), BG_LIGHT)
    tb(s3, c["stat"], cx + Inches(0.25), cy + ch - Inches(0.7), Inches(1.2), Inches(0.55),
       size=20, bold=True, color=c["color"])
    tb(s3, c["stat_label"], cx + Inches(1.5), cy + ch - Inches(0.65), cw - Inches(1.6), Inches(0.5),
       size=11, color=MUTED)


# ── SLIDE 4: Dashboard Features ─────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
bg(s4, BG_DARK)
box(s4, Inches(0), Inches(0), Inches(0.12), H, ORANGE)

tb(s4, "Dashboard 做了什麼？", Inches(0.5), Inches(0.4), Inches(12), Inches(0.7),
   size=30, bold=True, color=WHITE)
tb(s4, "github.io/KIWI — 手機可用，純靜態，零後端",
   Inches(0.5), Inches(1.05), Inches(12), Inches(0.45),
   size=15, color=MUTED, italic=True)
box(s4, Inches(0.5), Inches(1.55), Inches(4.5), Inches(0.05), ORANGE)

# feature grid 3×2
features = [
    ("📊  三合一 Gauge 環形圖",
     "CPI · TSI · AVI V5 各自帶動畫的環形儀表盤，一眼看出三個系統狀態",
     BLUE),
    ("⚡  Alert Banner",
     "TSI / CPI 達到警戒等級時，頂部自動出現彩色警告橫幅",
     ORANGE),
    ("🔢  指標拆解 Bar",
     "CPI 12 個 + TSI 9 個指標各自顯示進度條與 NORMAL / ELEVATED / HIGH 狀態",
     GREEN),
    ("🗓  H2 事件日曆",
     "NVIDIA 財報 / FOMC / 季節性弱勢月 — 關鍵日期一覽，搭配監控建議",
     YELLOW),
    ("📋  行動手冊",
     "DO / DON'T / TRIGGER 三欄：根據目前三系統狀態直接告訴你該做什麼",
     PURPLE),
    ("🌐  EN / 中文切換",
     "全頁雙語，localStorage 記憶設定，下次開啟自動套用",
     MUTED),
]

cols = 3
for i, (title, desc, col) in enumerate(features):
    row = i // cols
    ci  = i % cols
    fx  = Inches(0.5) + ci * Inches(4.25)
    fy  = Inches(1.85) + row * Inches(2.5)
    fw  = Inches(4.0)
    fh  = Inches(2.3)
    box(s4, fx, fy, fw, fh, BG_CARD)
    accent_bar(s4, fx, fy, fh, col)
    tb(s4, title, fx + Inches(0.2), fy + Inches(0.15), fw - Inches(0.3), Inches(0.5),
       size=14, bold=True, color=WHITE)
    tb(s4, desc,  fx + Inches(0.2), fy + Inches(0.7),  fw - Inches(0.3), Inches(1.4),
       size=12, color=MUTED, wrap=True)

# URL bar at bottom
box(s4, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4), BG_LIGHT)
tb(s4, "🌐   gutinganthony.github.io/KIWI",
   Inches(0.7), Inches(6.97), Inches(10), Inches(0.35),
   size=13, color=BLUE, bold=True)


# ── SLIDE 5: Claude's Contribution ──────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
bg(s5, BG_DARK)
box(s5, Inches(0), Inches(0), Inches(0.12), H, PURPLE)

tb(s5, "Claude 帶來了什麼？", Inches(0.5), Inches(0.35), Inches(12), Inches(0.7),
   size=30, bold=True, color=WHITE)
tb(s5, "一個 Session，從零到部署 — AI 協作開發的實際樣子",
   Inches(0.5), Inches(1.0), Inches(12), Inches(0.45),
   size=15, italic=True, color=MUTED)
box(s5, Inches(0.5), Inches(1.5), Inches(5), Inches(0.05), PURPLE)

# Timeline on left
timeline_items = [
    ("閱讀 Context",    "讀取 6 份跨 branch 的開發備忘，還原完整上下文",        BLUE),
    ("理解三個系統",    "從原始碼理解 CPI / TSI / AVI 的演算法與指標定義",      GREEN),
    ("設計 Dashboard",  "規劃 mobile-first 版面：3 Gauge + 拆解 + 日曆 + 手冊", ORANGE),
    ("撰寫完整 HTML",   "700+ 行單一 HTML 檔案，深色主題、SVG 動畫、i18n 系統",  YELLOW),
    ("設定 CI/CD",      "GitHub Actions workflow，push 到 main 自動部署",         PURPLE),
    ("開 PR + 上線",    "commit → push → 開 PR → GitHub Pages 部署完成",          WHITE),
]
for j, (step, desc, col) in enumerate(timeline_items):
    ty = Inches(1.85) + j * Inches(0.82)
    # dot + line
    box(s5, Inches(0.55), ty + Inches(0.12), Inches(0.18), Inches(0.18), col)
    if j < len(timeline_items) - 1:
        box(s5, Inches(0.62), ty + Inches(0.3), Inches(0.04), Inches(0.6), BG_LIGHT)
    tb(s5, step, Inches(0.9), ty, Inches(2.4), Inches(0.4),
       size=13, bold=True, color=col)
    tb(s5, desc, Inches(0.9), ty + Inches(0.38), Inches(5.5), Inches(0.38),
       size=11, color=MUTED)

# Right: key numbers
box(s5, Inches(7.2), Inches(1.75), Inches(5.7), Inches(5.4), BG_CARD)
tb(s5, "成果數字", Inches(7.5), Inches(1.9), Inches(5), Inches(0.5),
   size=14, bold=True, color=MUTED)
box(s5, Inches(7.4), Inches(2.38), Inches(5.3), Inches(0.04), BG_LIGHT)

stats = [
    ("1",     "個 Session 完成全部工作",  WHITE),
    ("67",    "個檔案 committed",          BLUE),
    ("3",     "個系統整合到一頁",          GREEN),
    ("100%",  "危機偵測率（CPI 4/4, TSI 7/7, AVI 5/5）", YELLOW),
    ("0",     "行後端程式碼（純靜態）",    ORANGE),
    ("∞",     "手機上的可及性",            PURPLE),
]
for k, (num, label, col) in enumerate(stats):
    sy = Inches(2.55) + k * Inches(0.72)
    box(s5, Inches(7.4), sy, Inches(5.3), Inches(0.65), BG_LIGHT if k % 2 == 0 else BG_CARD)
    tb(s5, num,   Inches(7.5), sy + Inches(0.08), Inches(1.0), Inches(0.5),
       size=22, bold=True, color=col)
    tb(s5, label, Inches(8.65), sy + Inches(0.16), Inches(3.8), Inches(0.4),
       size=12, color=MUTED)

# bottom quote
box(s5, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55), BG_CARD)
tb(s5,
   "「AVI / CPI / TSI 全部讀完後，Claude 設計架構、寫程式、開 PR、上線 — 我只需要說要做什麼。」",
   Inches(0.7), Inches(6.88), Inches(11.8), Inches(0.48),
   size=11, italic=True, color=MUTED)


# ── Save ──
out = "/home/user/KIWI/KIWI_Dashboard_Intro.pptx"
prs.save(out)
print(f"Saved: {out}")
