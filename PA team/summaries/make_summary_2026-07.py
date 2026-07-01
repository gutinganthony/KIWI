# -*- coding: utf-8 -*-
"""PA Team 四週市場回顧 (2026/06/02–07/01, 14 reports) on Cathay UB PB template.
   Report-style narrative prose + tables/cards. Output: 9 slides. All fonts 微軟正黑體."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SAMPLE = "/tmp/claude-0/-home-user-KIWI/9135307e-1a7d-5ee9-8555-29af972de8d1/scratchpad/PA_sample.pptx"
MEDIA  = "/tmp/claude-0/-home-user-KIWI/9135307e-1a7d-5ee9-8555-29af972de8d1/scratchpad/sample_x/ppt/media"
OUT    = "/home/user/KIWI/PA team/summaries/PA_Team_市場月報_2026-06-02_to_2026-07-01.pptx"

GOLD    = RGBColor(0xBE, 0x9E, 0x5A); GOLD_DK = RGBColor(0x97, 0x7B, 0x3B)
GOLD_LT = RGBColor(0xD9, 0xC7, 0x9F); NAVY    = RGBColor(0x23, 0x51, 0x73)
SLATE   = RGBColor(0x6A, 0x7C, 0x8D); INK     = RGBColor(0x22, 0x22, 0x22)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF); PAPER   = RGBColor(0xF6, 0xF3, 0xEC)
CARD    = RGBColor(0xFA, 0xF8, 0xF3); OWGREEN = RGBColor(0x3F, 0x7D, 0x54)
UWRED   = RGBColor(0xB1, 0x4A, 0x42)

E = Inches
W, H = E(13.333), E(7.5)
prs = Presentation(SAMPLE)

sldIdLst = prs.slides._sldIdLst
pres_part = prs.part
for sldId in list(sldIdLst):
    rId = sldId.get(qn('r:id'))
    try: pres_part.drop_rel(rId)
    except Exception: pass
    sldIdLst.remove(sldId)

BLANK   = prs.slide_layouts[7]
CONTENT = prs.slide_layouts[2]
FONT = "微軟正黑體"
_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def set_theme_font(face=FONT):
    from lxml import etree
    master = prs.slide_masters[0]; theme_part = None
    for rel in master.part.rels.values():
        if rel.reltype.endswith('/theme'): theme_part = rel.target_part; break
    if theme_part is None: return
    root = etree.fromstring(theme_part.blob); ns = {'a': _A}
    for tag in ('majorFont', 'minorFont'):
        fe = root.find(f'.//a:{tag}', ns)
        if fe is None: continue
        for sub in ('latin', 'ea', 'cs'):
            el = fe.find(f'a:{sub}', ns)
            if el is None: el = etree.SubElement(fe, qn(f'a:{sub}'))
            el.set('typeface', face)
        for f in fe.findall('a:font', ns):
            if f.get('script') in ('Hant', 'Hans'): f.set('typeface', face)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

def _force_run(r_el, face=FONT):
    rPr = r_el.find(qn('a:rPr'))
    if rPr is None: rPr = r_el.makeelement(qn('a:rPr'), {}); r_el.insert(0, rPr)
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None: el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', face)

def force_all_fonts(face=FONT):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for r in shape.text_frame._txBody.iter(qn('a:r')): _force_run(r, face)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for r in cell.text_frame._txBody.iter(qn('a:r')): _force_run(r, face)

set_theme_font()

def md(text):
    return [(part, i % 2 == 1) for i, part in enumerate(text.split("**")) if part]

def box(slide, x, y, w, h, fill, line=None, lw=0.75):
    sp = slide.shapes.add_shape(1, E(x), E(y), E(w), E(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else: sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def set_alpha(shape, opacity_pct):
    srgb = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))}))

def tbox(slide, x, y, w, h, anchor=None):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tb.text_frame.word_wrap = True
    if anchor: tb.text_frame.vertical_anchor = anchor
    return tb

def para(tf, segments, size, color, first=False, align=PP_ALIGN.LEFT,
         space_after=7, line=1.22, bold=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line
    for (t, b) in segments:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = (b or bold); r.font.color.rgb = color
    return p

def prose(slide, x, y, w, h, paragraphs, size=13, color=INK, space_after=8, line=1.25):
    tf = tbox(slide, x, y, w, h).text_frame
    for i, ptxt in enumerate(paragraphs):
        para(tf, md(ptxt), size, color, first=(i == 0), space_after=space_after, line=line)

def title_of(slide, text, size=20):
    ph = slide.shapes.title; ph.text = ""
    p = ph.text_frame.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK
    return ph

def drop_ph(slide, idxs):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in idxs:
            ph._element.getparent().remove(ph._element)

def source(slide, text):
    tf = tbox(slide, 0.3, 7.16, 10.5, 0.32).text_frame
    para(tf, md(text), 10.5, SLATE, first=True, bold=True, space_after=0)

def content_slide(heading, size=20):
    s = prs.slides.add_slide(CONTENT); drop_ph(s, {2})
    title_of(s, heading, size)
    box(s, 0.3, 1.12, 12.73, 0.045, GOLD)
    return s

def _no_style(tbl):
    tblPr = tbl._tbl.tblPr; tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    for c in list(tblPr):
        if c.tag == qn('a:tableStyleId'): tblPr.remove(c)
    el = tblPr.makeelement(qn('a:tableStyleId'), {}); el.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tblPr.append(el)

def _cell(c, text, size, bold, color, fill, align=PP_ALIGN.LEFT):
    c.margin_left = E(0.1); c.margin_right = E(0.08); c.margin_top = E(0.04); c.margin_bottom = E(0.04)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb = fill
    p = c.text_frame.paragraphs[0]; p.alignment = align; p.line_spacing = 1.1
    for (t, b) in md(text):
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = (b or bold); r.font.color.rgb = color

def table(slide, x, y, w, rows, widths, head_size=11.5, body_size=11, row_h=0.36):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, E(x), E(y), E(w), E(row_h * nr)); t = gt.table; _no_style(t)
    for j, wd in enumerate(widths): t.columns[j].width = E(wd)
    for i in range(nr): t.rows[i].height = E(row_h)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if i == 0: _cell(t.cell(i, j), val, head_size, True, WHITE, GOLD)
            else: _cell(t.cell(i, j), val, body_size, False, INK, PAPER if i % 2 == 1 else WHITE)
    return gt

def card(slide, x, y, w, h, header, header_color, items, body_size=11.5, item_color=INK, head_size=13):
    box(slide, x, y, w, h, CARD); box(slide, x, y, w, 0.5, header_color); box(slide, x, y, 0.07, h, header_color)
    tf = tbox(slide, x + 0.18, y + 0.05, w - 0.3, 0.42, MSO_ANCHOR.MIDDLE).text_frame
    para(tf, md(header), head_size, WHITE, first=True, bold=True, space_after=0)
    tf2 = tbox(slide, x + 0.18, y + 0.62, w - 0.34, h - 0.72).text_frame
    for i, it in enumerate(items):
        para(tf2, md("· " + it), body_size, item_color, first=(i == 0), space_after=6, line=1.16)

def chip(slide, x, y, w, h, num, label, col):
    box(slide, x, y, w, h, CARD); box(slide, x, y, 0.06, h, col)
    tf = tbox(slide, x + 0.16, y + 0.08, w - 0.24, 0.5).text_frame
    para(tf, md(num), 17, col, first=True, bold=True, space_after=0)
    tf2 = tbox(slide, x + 0.16, y + 0.56, w - 0.24, h - 0.6).text_frame
    para(tf2, md(label), 9.5, INK, first=True, space_after=0, line=1.05)


# ═══ SLIDE 1 — Title ═══
s = prs.slides.add_slide(BLANK)
for ph in list(s.placeholders): ph._element.getparent().remove(ph._element)
s.shapes.add_picture(f"{MEDIA}/image3.jpeg", 0, 0, W, H)
ov = box(s, 0, 0, 13.333, 7.5, RGBColor(0x12, 0x1C, 0x26)); set_alpha(ov, 46)
pn = box(s, 0, 3.2, 8.2, 2.25, RGBColor(0x0E, 0x18, 0x22)); set_alpha(pn, 40)
box(s, 0.32, 3.49, 6.4, 0.028, GOLD); box(s, 0.32, 5.27, 6.4, 0.028, GOLD)
tf = tbox(s, 0.34, 3.56, 9.5, 0.95).text_frame
p = tf.paragraphs[0]; p.line_spacing = 1.0
for t, sz in [("Portfolio Advisory", 31), ("  晨會重點摘要", 31)]:
    r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = WHITE
tf2 = tbox(s, 0.36, 4.5, 8.8, 0.6).text_frame
para(tf2, md("四週市場回顧：AI 盈餘牛市 × 鷹派轉向 × 月底大輪動"), 18, GOLD_LT, first=True, bold=True, space_after=0)
tf3 = tbox(s, 0.36, 5.4, 9.2, 0.5).text_frame
para(tf3, md("彙整期間 2026/06/02 – 07/01　·　14 份每日晨會重點摘要"), 13.5, WHITE, first=True, space_after=0)
s.shapes.add_picture(f"{MEDIA}/image1.png", E(9.07), E(0.33), E(3.8), E(0.6))
box(s, 0.44, 6.82, 0.42, 0.028, GOLD); box(s, 12.47, 6.82, 0.42, 0.028, GOLD)
tf4 = tbox(s, 0.32, 6.88, 5.0, 0.3).text_frame
para(tf4, md("私人銀行專屬投組諮詢科"), 11, GOLD_LT, first=True, bold=True, space_after=0)


# ═══ SLIDE 2 — Overview + timeline ═══
s = content_slide("總覽：四週市場主軸 — 盈餘牛市、鷹派轉向，月底掀起大輪動")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "過去四週（6/2–7/1），市場的核心張力可一句話概括：**基本面（企業盈餘）依然強勁，但價格的驅動力正在換手**。月初高盛、德銀、花旗接連上調標普 500 年底目標至 **8,000–8,100 點**；**6/5** 一場由部位過度擁擠引發的科技股重挫（那斯達克 **−4.18%**、費半 **−10.26%**），揭開「牛市中場、容錯率極低」的序幕。",
    "行情隨後圍繞**油價與聯儲**兩條主線拉鋸：**6/16** 美伊臨時協議使布倫特跌破 84；**6/24** 沃什主持的首場 FOMC 雖按兵不動，點陣圖卻意外轉鷹，市場將首次加息預期由 2027 年 3 月提前至 **2026 年 10 月**，華爾街集體撤回降息押注。",
    "到了**月底，劇情再度轉折**：一方面油價尾部風險明顯消退（波斯灣出口恢復、高盛下調 Q4 布倫特至 **80 美元**），宏觀轉趨平靜；另一方面市場內部掀起 **2003 年以來最大輪動**——資金自擁擠的大型科技與 AI「掏錢者」，輪入小盤股、景氣循環與 AI「收錢者」（記憶體），同時 BIS 對 AI 資本支出泡沫拉響警報。**焦點正從宏觀（油價、利率）轉向微觀（AI 估值與盈餘能否兌現）**。",
], size=12, space_after=8)
tlx = 8.1
box(s, tlx, 1.35, 4.95, 5.62, CARD); box(s, tlx, 1.35, 4.95, 0.44, NAVY)
tf = tbox(s, tlx + 0.2, 1.38, 4.6, 0.4, MSO_ANCHOR.MIDDLE).text_frame
para(tf, md("四週關鍵節點"), 13, WHITE, first=True, bold=True, space_after=0)
timeline = [
    ("6/02–04", "AI 熱撞上油價衝擊；GS/DB 上調標普目標至 8000"),
    ("6/05", "科技重挫：那指 −4.18%、費半 −10.26%，韓股熔斷"),
    ("6/09–11", "「健康調整」非見頂；大摩流動性指標轉緊"),
    ("6/16–17", "美伊臨時協議，布倫特跌破 84；大型科技去擁擠"),
    ("6/18", "擁擠度創高：做多半導體＝史上最擁擠（80%）"),
    ("6/24", "FOMC 鷹派轉向：首次加息預期提前至 2026/10"),
    ("6/25", "大摩資產配置：美股超配、偏好中期美債"),
    ("6/30", "高盛亞股 H2：超配北亞，KOSPI 12000、台股 51000"),
    ("7/01", "BIS 泡沫警報＋大輪動：掏錢者→收錢者、油價尾部消退"),
]
ty = 1.9
for date, txt in timeline:
    tf = tbox(s, tlx + 0.2, ty, 1.12, 0.5).text_frame
    para(tf, md(date), 10.5, GOLD_DK, first=True, bold=True, space_after=0)
    tf2 = tbox(s, tlx + 1.28, ty, 3.5, 0.6).text_frame
    para(tf2, md(txt), 9.5, INK, first=True, space_after=0, line=1.05)
    ty += 0.558
source(s, "資料來源：Goldman Sachs / Morgan Stanley / Deutsche Bank / Citi / BofA / BIS / Bloomberg（PA 晨會摘要 2026/06–07）")


# ═══ SLIDE 3 — AI earnings bull ═══
s = content_slide("主題一：AI 盈餘牛市 — 盈餘紮實，但「掏錢者」開始出血")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "本輪美股的核心驅動力，是一場被花旗稱為「史無前例」的 AI 資本開支超級週期，且有紮實盈餘背書：標普 500 第一季實際盈餘較市場共識高約 **13%**，2026 年 EPS 自 320 上修至 **340–350 美元**、2027 上看 **385–400**；第二季盈餘增速預期更升至 **22%**，為 2021 年以來季報前最高。高盛、花旗、大摩目標分別為 **8000、8100、8300 點**。",
    "但獲利高度集中、且產業鏈內部開始分化。中位數公司盈餘增速僅 **9%**，AI 基礎設施貢獻標普第二季盈餘增量近 **60%**，輝達與美光合計逾 **40%**。更關鍵的是，6 月市場開始「投票」：AI 資本支出的「掏錢者」（超大規模雲端）單月重挫 **18%**（Meta 上市以來最差），「收錢者」（記憶體籃子）卻逆勢 **+16%**。",
    "支撐「收錢者」的是記憶體超級週期：DRAM/NAND 今年創紀錄供給缺口，2027 年可能擴大、2028 年仍深度為負，為記憶體與北亞供應鏈提供多年能見度。這也解釋了為何配置上宜「重收錢者、輕掏錢者」。",
], size=12, space_after=8)
table(s, 8.1, 1.5, 4.95, [
    ["指標", "現況"],
    ["標普目標", "GS 8000・Citi 8100・MS 8300"],
    ["2026/27 EPS", "340–350 / 385–400 美元"],
    ["Q2 盈餘增速", "22%（中位數僅 9%）"],
    ["AI 貢獻", "佔標普 Q2 盈餘增量 ~60%"],
    ["集中度", "輝達＋美光 > 40%"],
    ["6 月分化", "掏錢者 −18% / 收錢者 +16%"],
], widths=[1.35, 3.6], row_h=0.62, body_size=10.5)
source(s, "資料來源：Goldman Sachs / Citi / Morgan Stanley / BIS（PA 晨會摘要 2026/06–07）")


# ═══ SLIDE 4 — inflation + Fed pivot ═══
s = content_slide("主題二：通膨回溫與聯儲鷹派轉向 — 1996 還是 1999？")
prose(s, 0.3, 1.35, 12.7, 3.0, [
    "這段期間最關鍵的宏觀轉折，是**通膨預期回溫與聯儲態度轉鷹**。德銀以「1999 遇上 1990」形容當前處境——AI 熱潮像 1999、油價衝擊像 1990；其預測中增長僅小修、通膨卻大幅上修，迫使央行週期轉向。轉折點落在 **6/24**：沃什主持的首場 FOMC 維持利率不變，但點陣圖意外轉鷹（**18 位官員中 9 位**指向年內加息），並刪除「寬鬆偏向」措辭。掉期市場隨即將首次加息預期由 2027 年 3 月提前至 **2026 年 10 月**，年內已定價約 **41 個基點**。聯邦基金利率已在 **3.5–4.5%** 區間維持逾 18 個月，市場預期至少再維持一年。",
    "市場的核心辯論，是這輪 AI 繁榮究竟是 **1996（生產率紅利，可從容按兵不動）還是 1999（需求過熱，須預先收緊）**。沃什明顯傾向 1996，但面對關稅、財政赤字擴張與全球化紅利消退，其通膨壓力大於格林斯潘當年——這使各大投行對聯儲路徑出現罕見分歧（惟月底油價回落已略微削弱加息理由）：",
], size=12, space_after=9)
table(s, 0.3, 4.72, 12.73, [
    ["機構", "對聯儲路徑的看法"],
    ["美銀 BofA", "年內加息 3 碼（9/10/12 月，共 +75bp）；維持「逢高減倉」"],
    ["德銀 DB", "年內加息 2 碼至 4.1%，警告 7 月可能提前"],
    ["巴克萊 Barclays", "建議做空美債，殖利率曲線「熊陡」，各天期目標上調約 35bp"],
    ["花旗 Citi", "逆勢，仍看降息（10/12 月、明年 1 月各 −25bp）"],
    ["大摩 MS", "今年按兵不動，明年第一季降息 1 碼"],
], widths=[2.3, 10.43], row_h=0.38, body_size=11)
source(s, "資料來源：BofA / Deutsche Bank / Barclays / Citi / Morgan Stanley（PA 晨會摘要 2026/06）")


# ═══ SLIDE 5 — oil tail receding + rates ═══
s = content_slide("主題三：油價尾部消退，利率仍是紅線 —「宏觀平靜、微觀喧囂」")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "油價曾是這輪行情的「總開關」：伊朗戰爭下，德銀設定基準布倫特 **86**、事故情境 **150**。但 **6/16** 美伊臨時協議後情勢逐步緩解——月底波斯灣石油出口已恢復至正常水準的 **66%**，高盛將第四季布倫特預測下調至 **80 美元**，10 年期盈虧平衡通膨降至一年多低點、G10 外匯波動更跌至五年低位。**油價這個最大的宏觀尾部風險，正明顯消退**。",
    "利率則仍是紅線：高盛視 **10 年期美債 5%** 為股市估值的臨界，目前尚未觸及；聯邦基金利率預期在 3.5–4.5% 續守，但未來 2–3 個月聯儲鷹派風險仍會擾動資產。財政赤字近 GDP **6.6%**、長端供給壓力使曲線維持熊平。",
    "宏觀噪音下降的結果，是市場進入**「宏觀平靜、微觀喧囂」**格局：下檔風險的主角，已從油價與 CPI，換成 AI 估值、盈餘持續性與資本支出回報。",
], size=12, space_after=8)
card(s, 8.1, 1.5, 4.95, 1.5, "油價：尾部消退", GOLD_DK,
     ["波斯灣出口回升至正常 66%", "GS 下調 Q4 布倫特至 80（基準曾 86 / 尾部 150）"], body_size=11)
card(s, 8.1, 3.18, 4.95, 1.5, "利率：仍是紅線", NAVY,
     ["10 年期美債 5% ＝股市紅線", "Fed funds 3.5–4.5% 續守；曲線熊平"], body_size=11)
card(s, 8.1, 4.86, 4.95, 1.5, "格局轉變", UWRED,
     ["宏觀平靜、微觀喧囂", "焦點：油價 CPI → AI 估值與盈餘"], body_size=11)
source(s, "資料來源：Deutsche Bank / Goldman Sachs / Bloomberg（PA 晨會摘要 2026/06–07）")


# ═══ SLIDE 6 — The Great Rotation ═══
s = content_slide("主題四：2003 年以來最大輪動 —「掏錢者→收錢者、大盤→小盤」")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "六月底至七月初，美股掀起 **2003 年以來最大的結構性輪動**。資金自擁擠的大型科技／半導體與 AI「掏錢者」（超大規模雲端），輪入**小盤股、生技、航空旅遊、住房、REITs 與區域銀行**，以及 AI 產業鏈的「收錢者」（記憶體）。羅素 2000 對標普 500 的年初至今超額已達約 **1,240 個基點**——若延續至年末，將是 2003 年以來最大的年度超額。",
    "背後是對 AI 資本支出的重新定價。全球五大超大規模雲端 2025–26 年資本支出合計逾 **1 兆美元**，BIS 因而拉響警報：對照 1840 年代鐵路熱與 1990 年代末網路泡沫，均源於真實技術突破、卻吸引超出回報所能支撐的資本。安聯亦將 SpaceX 完成 860 億美元 IPO 後再啟 250 億美元債券發行，視為進入「泡沫區域」的訊號。",
    "這並非泡沫破裂的定論，而是**機率分布的位移**——上行驚喜門檻升高、下行尾部觸發條件降低。倉位顯示：超大規模雲端多頭已降至近三年低點、AI 基礎設施多頭處近三年高位，對沖基金淨賣美股速度創「關稅日」以來最快，但淨敞口仍高居第 **87 百分位**，部位並不輕盈。對策：**續抱股票多頭、同時以期權／做多波動率保護下行**。",
], size=12, space_after=8)
chips = [
    ("+1,240bp", "羅素 2000 vs 標普 YTD 超額", NAVY),
    ("−18% / +16%", "6 月 掏錢者 / 收錢者", UWRED),
    ("> $1 兆", "五大雲端 25–26 資本支出", GOLD_DK),
    ("87th", "對沖基金淨敞口百分位（仍高）", NAVY),
]
cy = 1.5
for num, lab, col in chips:
    chip(s, 8.1, cy, 4.95, 1.2, num, lab, col); cy += 1.32
source(s, "資料來源：Goldman Sachs / BIS / Allianz / Bloomberg（PA 晨會摘要 2026/07/01）")


# ═══ SLIDE 7 — Regional + temperature ═══
s = content_slide("區域分化與市場溫度")
prose(s, 0.3, 1.32, 12.7, 1.65, [
    "各區域分化明顯：**美國**最具韌性但廣度外擴；**歐洲**最弱、近技術性衰退；**亞洲（除日）為最大亮點**——上半年 MXAPJ **+21%**，韓國 **+119%**、台灣 **+56%** 領漲，高盛下半年「超配北亞」並給出 KOSPI **12,000**、台股 **51,000** 的目標；**日本**受油價與匯率影響、**中國**資金面偏弱。與此同時，情緒與擁擠度在月中觸頂後開始去化，但淨部位仍不輕盈。",
], size=12, space_after=6)
table(s, 0.3, 3.15, 7.45, [
    ["區域", "觀點", "關鍵數據"],
    ["美國", "韌性，廣度外擴", "Q2 盈餘 22%（中位數 9%）"],
    ["亞洲(除日)", "最大亮點，超配北亞", "H1 MXAPJ +21%；韓+119%/台+56%"],
    ["韓國/台灣", "GS 明確超配", "KOSPI 12000、台股 51000 目標"],
    ["歐洲", "最弱", "2026 GDP 0.5%；近技術性衰退"],
    ["亞洲低配", "GS：澳/泰/印尼/菲", "中國 A 超配；星港印中性"],
], widths=[1.25, 2.35, 3.85], row_h=0.5, body_size=10.5, head_size=11)
stats = [
    ("39%", "AI 股佔標普權重", GOLD_DK),
    ("22%", "Q2 盈餘增速(中位僅9%)", NAVY),
    ("9.2", "美銀牛熊(賣出區)", UWRED),
    ("4.1%", "平均現金水位", NAVY),
    ("87th", "對沖基金淨敞口(仍高)", UWRED),
    ("28%", "視 AI 泡沫為首要尾部風險", GOLD_DK),
]
bx, by, cw, chh = 7.95, 3.15, 2.5, 1.12
for i, (num, lab, col) in enumerate(stats):
    chip(s, bx + (i % 2) * (cw + 0.08), by + (i // 2) * (chh + 0.1), cw, chh, num, lab, col)
source(s, "資料來源：Goldman Sachs / Morgan Stanley / BofA / Bloomberg（PA 晨會摘要 2026/06–07）")


# ═══ SLIDE 8 — Conclusion + allocation ═══
s = content_slide("結論與資產配置建議")
prose(s, 0.3, 1.32, 12.7, 1.45, [
    "綜合四週觀點，基調由「建設性、但防禦微調」進一步演化為**「續抱多頭、但加速廣度輪動並強化下行保護」**：油價尾部消退令宏觀轉平靜，但 AI 資本支出的回報與泡沫疑慮成為新的核心風險。宜順勢參與輪動、降低對擁擠「掏錢者」的集中、並以波動率／期權對沖。具體配置：",
], size=12, space_after=6)
cy, ch2 = 2.78, 3.8
card(s, 0.3, cy, 4.05, ch2, "▲  增持 OVERWEIGHT", OWGREEN, [
    "北亞：韓國・台灣・日本（＋中國 A）— GS 超配、晶片超級週期",
    "AI「收錢者」記憶體 ＞「掏錢者」超大規模雲端",
    "輪動受益：小盤股、生技、景氣循環（航空/旅遊/REITs/區域銀行）",
    "中期美債；美股金融・工業・醫療（廣度外擴）",
], body_size=11)
card(s, 4.64, cy, 4.05, ch2, "▼  減持 UNDERWEIGHT", UWRED, [
    "擁擠的超大規模雲端「掏錢者」/ 大盤 AI 集中部位",
    "長天期債券：熊陡，10Y 逼近 5%",
    "亞洲低配：澳洲・泰國・印尼・菲律賓（GS）",
    "歐洲股票 / 信用；日圓 / 未對沖日股",
], body_size=11)
card(s, 8.98, cy, 4.05, ch2, "●  中性 / 戰術 NEUTRAL", GOLD_DK, [
    "波動率 / 期權：BIS 泡沫警報下的下行保護（升級）",
    "黃金：對沖；現金：~4% 緩衝",
    "離岸中國 / 印度 / 星港：待盈利改善（GS 中性）",
    "歐洲能源股：Hormuz 重啟 → 中性",
], body_size=11)
tf = tbox(s, 0.3, 6.72, 12.7, 0.38).text_frame
para(tf, md("※ 上述為彙整各投行研究之整合視角，非個別投資建議；花旗 / 大摩對聯儲路徑（仍見降息）與多數機構（轉向加息）存在分歧，請依自身風險屬性判讀。"),
     9, SLATE, first=True, space_after=0)


# ═══ SLIDE 9 — Watch list ═══
s = content_slide("下一步：關鍵觀察點與風險地圖")
prose(s, 0.3, 1.32, 12.7, 1.0, [
    "隨油價尾部消退，焦點轉向 AI 微觀與輪動能否延續。以下六個變數將決定劇本走向（1996 軟著陸 vs 1999 過熱），建議作為動態調整的觸發點：",
], size=12, space_after=4)
watch = [
    ("Q2 財報（7 月底）", "AI 資本支出能否轉化為實質營收——輝達/美光/雲端指引是關鍵", GOLD_DK),
    ("融資條件（BIS）", "若回報不及預期，融資可能驟緊、令 capex 繁榮急轉直下", UWRED),
    ("大輪動的持續性", "小盤/循環 vs 大盤科技：是新主線，還是 AI 引擎的短暫喘息", NAVY),
    ("10 年期美債 / Fed", "5% 紅線；Fed 3.5–4.5% 續守，未來 2–3 月鷹派風險", NAVY),
    ("東西成本之爭", "GLM-5.2 全程未用輝達晶片，挑戰「capex 須持續擴張」假設", GOLD_DK),
    ("集中度風險", "輝達＋美光 >40% 貢獻；「掏錢者」若減支出即為斷裂點", UWRED),
]
gx, gy, gw, gh = 0.3, 2.5, 4.13, 1.78
for i, (t, d, col) in enumerate(watch):
    x = gx + (i % 3) * (gw + 0.07); y = gy + (i // 3) * (gh + 0.12)
    box(s, x, y, gw, gh, CARD); box(s, x, y, gw, 0.07, col); box(s, x, y, 0.07, gh, col)
    tf = tbox(s, x + 0.2, y + 0.14, gw - 0.34, 0.5).text_frame
    para(tf, md(t), 13, col, first=True, bold=True, space_after=0)
    tf2 = tbox(s, x + 0.2, y + 0.62, gw - 0.36, gh - 0.74).text_frame
    para(tf2, md(d), 10.5, INK, first=True, space_after=0, line=1.16)
box(s, 0.3, 6.62, 12.73, 0.5, NAVY)
tf = tbox(s, 0.45, 6.66, 12.4, 0.42, MSO_ANCHOR.MIDDLE).text_frame
para(tf, md("底線：宏觀轉平靜、微觀更喧囂 —— 參與輪動、分散集中、以波動率對沖，緊盯 AI 資本支出的回報兌現。"),
     12, WHITE, first=True, bold=True, space_after=0)


force_all_fonts()
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
