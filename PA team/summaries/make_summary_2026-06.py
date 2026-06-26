# -*- coding: utf-8 -*-
"""PA Team 四週市場回顧 — rebuilt on the user's Cathay UB Private Banking template.
   Report-style narrative prose + tables/cards. Output: 8 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SAMPLE = "/tmp/claude-0/-home-user-KIWI/9135307e-1a7d-5ee9-8555-29af972de8d1/scratchpad/PA_sample.pptx"
MEDIA  = "/tmp/claude-0/-home-user-KIWI/9135307e-1a7d-5ee9-8555-29af972de8d1/scratchpad/sample_x/ppt/media"
OUT    = "/home/user/KIWI/PA team/summaries/PA_Team_市場月報_2026-06-02_to_2026-06-25.pptx"

# ── Palette (from template theme) ──
GOLD    = RGBColor(0xBE, 0x9E, 0x5A)   # accent4
GOLD_DK = RGBColor(0x97, 0x7B, 0x3B)   # accent2
GOLD_LT = RGBColor(0xD9, 0xC7, 0x9F)   # accent3
NAVY    = RGBColor(0x23, 0x51, 0x73)   # dk2
SLATE   = RGBColor(0x6A, 0x7C, 0x8D)   # lt2
INK     = RGBColor(0x22, 0x22, 0x22)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PAPER   = RGBColor(0xF6, 0xF3, 0xEC)
CARD    = RGBColor(0xFA, 0xF8, 0xF3)
OWGREEN = RGBColor(0x3F, 0x7D, 0x54)
UWRED   = RGBColor(0xB1, 0x4A, 0x42)

E = Inches
W, H = E(13.333), E(7.5)

prs = Presentation(SAMPLE)

# ── remove the 6 sample slides (drop parts + rels so no duplicates) ──
sldIdLst = prs.slides._sldIdLst
pres_part = prs.part
for sldId in list(sldIdLst):
    rId = sldId.get(qn('r:id'))
    try:
        pres_part.drop_rel(rId)
    except Exception:
        pass
    sldIdLst.remove(sldId)

BLANK   = prs.slide_layouts[7]   # 空白
CONTENT = prs.slide_layouts[2]   # 自訂版面配置 (logo + watermark + line + page no.)

FONT = "微軟正黑體"
_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def set_theme_font(face=FONT):
    """Rewrite theme major/minor fonts so ALL inherited text uses `face`."""
    from lxml import etree
    master = prs.slide_masters[0]
    theme_part = None
    for rel in master.part.rels.values():
        if rel.reltype.endswith('/theme'):
            theme_part = rel.target_part; break
    if theme_part is None:
        return
    root = etree.fromstring(theme_part.blob)
    ns = {'a': _A}
    for tag in ('majorFont', 'minorFont'):
        fe = root.find(f'.//a:{tag}', ns)
        if fe is None:
            continue
        for sub in ('latin', 'ea', 'cs'):
            el = fe.find(f'a:{sub}', ns)
            if el is None:
                el = etree.SubElement(fe, qn(f'a:{sub}'))
            el.set('typeface', face)
        for f in fe.findall('a:font', ns):
            if f.get('script') in ('Hant', 'Hans'):
                f.set('typeface', face)
    theme_part._blob = etree.tostring(
        root, xml_declaration=True, encoding='UTF-8', standalone=True)

def _force_run(r_el, face=FONT):
    rPr = r_el.find(qn('a:rPr'))
    if rPr is None:
        rPr = r_el.makeelement(qn('a:rPr'), {}); r_el.insert(0, rPr)
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', face)

def force_all_fonts(face=FONT):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for r in shape.text_frame._txBody.iter(qn('a:r')):
                    _force_run(r, face)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for r in cell.text_frame._txBody.iter(qn('a:r')):
                            _force_run(r, face)

set_theme_font()


# ── helpers ─────────────────────────────────────────────────────
def md(text):
    """split **bold** markup -> [(text, bold), ...]"""
    return [(part, i % 2 == 1) for i, part in enumerate(text.split("**")) if part]

def box(slide, x, y, w, h, fill, line=None, lw=0.75):
    sp = slide.shapes.add_shape(1, E(x), E(y), E(w), E(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else:    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def set_alpha(shape, opacity_pct):
    srgb = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))}))

def tbox(slide, x, y, w, h, anchor=None):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tb.text_frame.word_wrap = True
    if anchor: tb.text_frame.vertical_anchor = anchor
    return tb

def para(tf, segments, size, color, first=False, align=PP_ALIGN.LEFT,
         space_after=7, line=1.22, bold=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for (t, b) in segments:
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = (b or bold); r.font.color.rgb = color
    return p

def prose(slide, x, y, w, h, paragraphs, size=13, color=INK, space_after=8, line=1.25):
    tf = tbox(slide, x, y, w, h).text_frame
    for i, ptxt in enumerate(paragraphs):
        para(tf, md(ptxt), size, color, first=(i == 0), space_after=space_after, line=line)

def title_of(slide, text, size=20):
    ph = slide.shapes.title
    ph.text = ""
    p = ph.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK
    return ph

def drop_ph(slide, idxs):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in idxs:
            ph._element.getparent().remove(ph._element)

def source(slide, text):
    tf = tbox(slide, 0.3, 7.16, 9.4, 0.32).text_frame
    para(tf, md(text), 10.5, SLATE, first=True, bold=True, space_after=0)

def content_slide(heading, size=20):
    s = prs.slides.add_slide(CONTENT)
    drop_ph(s, {2})                       # remove body prompt placeholder
    title_of(s, heading, size)
    box(s, 0.3, 1.12, 12.73, 0.045, GOLD) # gold rule under title
    return s

# ── No-style native table ──
def _no_style(tbl):
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    for c in list(tblPr):
        if c.tag == qn('a:tableStyleId'): tblPr.remove(c)
    el = tblPr.makeelement(qn('a:tableStyleId'), {})
    el.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tblPr.append(el)

def _cell(c, text, size, bold, color, fill, align=PP_ALIGN.LEFT):
    c.margin_left = E(0.1); c.margin_right = E(0.08)
    c.margin_top = E(0.04); c.margin_bottom = E(0.04)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb = fill
    p = c.text_frame.paragraphs[0]; p.alignment = align; p.line_spacing = 1.1
    for (t, b) in md(text):
        r = p.add_run(); r.text = t
        r.font.size = Pt(size); r.font.bold = (b or bold); r.font.color.rgb = color

def table(slide, x, y, w, rows, widths, head_size=11.5, body_size=11, row_h=0.36):
    nr, nc = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nr, nc, E(x), E(y), E(w), E(row_h * nr))
    t = gt.table; _no_style(t)
    for j, wd in enumerate(widths): t.columns[j].width = E(wd)
    for i in range(nr): t.rows[i].height = E(row_h)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if i == 0:
                _cell(t.cell(i, j), val, head_size, True, WHITE, GOLD,
                      align=PP_ALIGN.LEFT)
            else:
                fill = PAPER if i % 2 == 1 else WHITE
                _cell(t.cell(i, j), val, body_size, False, INK, fill,
                      align=PP_ALIGN.LEFT)
    return gt

def card(slide, x, y, w, h, header, header_color, items, body_size=11.5,
         item_color=INK, head_size=13):
    box(slide, x, y, w, h, CARD)
    box(slide, x, y, w, 0.5, header_color)             # header band
    box(slide, x, y, 0.07, h, header_color)            # left spine
    tf = tbox(slide, x + 0.18, y + 0.05, w - 0.3, 0.42, MSO_ANCHOR.MIDDLE).text_frame
    para(tf, md(header), head_size, WHITE, first=True, bold=True, space_after=0)
    tf2 = tbox(slide, x + 0.18, y + 0.62, w - 0.34, h - 0.72).text_frame
    for i, it in enumerate(items):
        para(tf2, md("· " + it), body_size, item_color, first=(i == 0),
             space_after=6, line=1.16)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (replicate template cover)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
for ph in list(s.placeholders):
    ph._element.getparent().remove(ph._element)
s.shapes.add_picture(f"{MEDIA}/image3.jpeg", 0, 0, W, H)         # skyline (full quality)
ov = box(s, 0, 0, 13.333, 7.5, RGBColor(0x12, 0x1C, 0x26)); set_alpha(ov, 46)
panel = box(s, 0, 3.2, 8.2, 2.25, RGBColor(0x0E, 0x18, 0x22)); set_alpha(panel, 40)
box(s, 0.32, 3.49, 6.4, 0.028, GOLD)                            # upper gold line
box(s, 0.32, 5.27, 6.4, 0.028, GOLD)                            # lower gold line
tf = tbox(s, 0.34, 3.56, 9.5, 0.95).text_frame
p = tf.paragraphs[0]; p.line_spacing = 1.0
for t, c, sz in [("Portfolio Advisory", WHITE, 31), ("  晨會重點摘要", WHITE, 31)]:
    r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = c
tf2 = tbox(s, 0.36, 4.5, 8.6, 0.6).text_frame
para(tf2, md("四週市場回顧：AI 盈餘牛市 × 油價衝擊 × 聯儲鷹派轉向"),
     18, GOLD_LT, first=True, bold=True, space_after=0)
tf3 = tbox(s, 0.36, 5.4, 9.0, 0.5).text_frame
para(tf3, md("彙整期間 2026/06/02 – 06/25　·　12 份每日晨會重點摘要"),
     13.5, WHITE, first=True, space_after=0)
s.shapes.add_picture(f"{MEDIA}/image1.png", E(9.07), E(0.33), E(3.8), E(0.6))  # logo
box(s, 0.44, 6.82, 0.42, 0.028, GOLD)
box(s, 12.47, 6.82, 0.42, 0.028, GOLD)
tf4 = tbox(s, 0.32, 6.88, 5.0, 0.3).text_frame
para(tf4, md("私人銀行專屬投組諮詢科"), 11, GOLD_LT, first=True, bold=True, space_after=0)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview + timeline
# ════════════════════════════════════════════════════════════════
s = content_slide("總覽：四週市場主軸 — 盈餘牛市撞上鷹派轉向")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "過去四週（6/2–6/25），市場的核心張力可以一句話概括：**基本面（企業盈餘）依然強勁，但宏觀環境（通膨與利率）正在轉向**。月初高盛、德銀與花旗接連上調標普 500 年底目標至 **8,000–8,100 點**，AI 資本開支超級週期帶動的盈餘上修是主要動力；然而 **6/5** 一場由部位過度擁擠引發的科技股重挫（那斯達克單日 **−4.18%**、費半 **−10.26%**），揭開了「牛市中場、容錯率極低」的序幕。",
    "隨後行情圍繞兩條主線拉鋸。其一是**油價與地緣**：伊朗戰爭推升通膨，布倫特原油在「基準 86、事故情境 150 美元」之間擺盪，直到 **6/16** 美伊達成臨時協議、霍爾木茲海峽重啟在望，油價才回落至 84 美元以下。其二是**聯儲政策**：**6/24** 沃什（Warsh）主持的首場 FOMC 雖按兵不動，點陣圖卻意外轉鷹，市場將首次加息預期由 2027 年 3 月大幅提前至 **2026 年 10 月**，華爾街集體撤回降息押注。",
    "換言之，這既非單純多頭、也非空頭，而是一個**「相信 AI 敘事」與「預期透支」相互拉扯**的市場：盈餘的確定性仍支撐估值，但擁擠度創歷史新高、利率逼近紅線，使任何利空都可能被槓桿與波動率放大。",
], size=12.5, space_after=9)
# right timeline card
tlx = 8.1
box(s, tlx, 1.35, 4.95, 5.62, CARD)
box(s, tlx, 1.35, 4.95, 0.46, NAVY)
tf = tbox(s, tlx + 0.2, 1.39, 4.6, 0.4, MSO_ANCHOR.MIDDLE).text_frame
para(tf, md("四週關鍵節點"), 13, WHITE, first=True, bold=True, space_after=0)
timeline = [
    ("6/02–04", "「1999 遇上 1990」：AI 熱潮撞上油價衝擊，GS/DB 上調標普目標至 8000"),
    ("6/05",    "科技重挫：那指 −4.18%、費半 −10.26%，逾兆美元蒸發，韓股熔斷"),
    ("6/09–11", "「健康調整」非見頂；惟大摩流動性指標轉緊"),
    ("6/16–17", "美伊臨時協議，布倫特跌破 84；大型科技去擁擠回中性"),
    ("6/18",    "擁擠度創高：做多半導體＝史上最擁擠（80%），現金升至 4.1%"),
    ("6/24",    "FOMC 鷹派轉向：首次加息預期提前至 2026/10"),
    ("6/25",    "資產配置防禦微調：美股超配、歐能源降中性、偏好中期美債"),
]
ty = 1.98
for date, txt in timeline:
    tf = tbox(s, tlx + 0.2, ty, 1.15, 0.5).text_frame
    para(tf, md(date), 11, GOLD_DK, first=True, bold=True, space_after=0)
    tf2 = tbox(s, tlx + 1.32, ty, 3.5, 0.7).text_frame
    para(tf2, md(txt), 10, INK, first=True, space_after=0, line=1.08)
    ty += 0.705
source(s, "資料來源：Goldman Sachs / Morgan Stanley / Deutsche Bank / Citi / BofA / Bloomberg（PA 晨會摘要 2026/06）")


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Theme 1: AI earnings bull
# ════════════════════════════════════════════════════════════════
s = content_slide("主題一：AI 盈餘牛市 — 中場，但窗口正在收窄")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "本輪美股的核心驅動力，是一場被花旗形容為「史無前例」的 AI 資本開支超級週期。與 1999 年或 2021 年的純情緒行情不同，這次有紮實的盈餘背書：標普 500 第一季實際盈餘較市場共識高出約 **13%**，2026 年每股盈餘預估自 320 美元一路上修至 **340–350 美元**，2027 年上看 **385–400 美元**。高盛、花旗因而分別將年底目標上調至 **8,000、8,100 點**，大摩更給出 2027 年中 **8,300 點**目標。",
    "但多家投行同時示警「窗口正在收窄」。後續上漲須靠盈餘、而非估值擴張推動——新目標隱含的本益比（約 20–23 倍）反而較先前更低；最快的盈餘增速可能已過，正向驚喜將逐步收斂，形成一個**容錯率極低**的市場。",
    "更值得警惕的是集中度：AI 相關個股已佔標普 500 權重達 **39%**。當「賣鏟子」（AI 基礎設施）邏輯被市場充分認識後，下行風險呈現不對稱放大——越是眾所周知的主題，一旦增速放緩，定價逆轉就越快。這正是後續配置須主動分散、降低集中的根本原因。",
], size=12.5, space_after=9)
table(s, 8.1, 1.55, 4.95, [
    ["指標", "現況"],
    ["標普目標", "GS 8000・Citi 8100・DB 8000"],
    ["（大摩）", "MS 8300（2027 年中）"],
    ["2026 EPS", "340–350 美元（↑自 320）"],
    ["2027 EPS", "385–400 美元"],
    ["Q1 盈餘", "較市場共識高約 13%"],
    ["AI 權重", "佔標普 500 約 39%"],
    ["本益比", "約 20–23 倍（隨盈餘上修而降）"],
], widths=[1.15, 3.8], row_h=0.55, body_size=10.5)
source(s, "資料來源：Goldman Sachs / Citi / Morgan Stanley（PA 晨會摘要 2026/06）")


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Theme 2: inflation + Fed hawkish pivot
# ════════════════════════════════════════════════════════════════
s = content_slide("主題二：通膨回溫與聯儲鷹派轉向 — 1996 還是 1999？")
prose(s, 0.3, 1.35, 12.7, 3.0, [
    "這段期間最關鍵的宏觀轉折，是**通膨預期回溫與聯儲態度轉鷹**。德銀以「1999 遇上 1990」形容當前處境——AI 熱潮像 1999、油價衝擊像 1990；其最新預測中增長僅小幅下修、通膨卻大幅上修，名義 GDP 被通膨托高，迫使央行週期轉向。轉折點落在 **6/24**：沃什主持的首場 FOMC 維持利率不變，但點陣圖意外轉鷹（**18 位官員中 9 位**指向年內加息），並刪除「寬鬆偏向」措辭、拒絕提供前瞻指引。掉期市場隨即將首次加息預期由 2027 年 3 月提前至 **2026 年 10 月**，年內已定價約 **41 個基點**，2 年期美債殖利率創 3 月以來最大單日漲幅。",
    "市場的核心辯論，是這輪 AI 繁榮究竟是 **1996（生產率紅利，可從容按兵不動）還是 1999（需求過熱，須預先收緊）**。沃什明顯傾向 1996 劇本，但面對關稅、財政赤字擴張與全球化紅利消退，其通膨壓力大於格林斯潘當年——這也使各大投行對聯儲路徑出現罕見分歧：",
], size=12.5, space_after=9)
table(s, 0.3, 4.62, 12.73, [
    ["機構", "對聯儲路徑的看法"],
    ["美銀 BofA", "年內加息 3 碼（9/10/12 月，共 +75bp）；維持「逢高減倉」"],
    ["德銀 DB", "年內加息 2 碼至 4.1%，警告 7 月可能提前"],
    ["巴克萊 Barclays", "建議做空美債，殖利率曲線「熊陡」，各天期目標上調約 35bp"],
    ["花旗 Citi", "逆勢，仍看降息（10/12 月、明年 1 月各 −25bp）"],
    ["大摩 MS", "今年按兵不動，明年第一季降息 1 碼"],
], widths=[2.3, 10.43], row_h=0.4, body_size=11)
source(s, "資料來源：BofA / Deutsche Bank / Barclays / Citi / Morgan Stanley（PA 晨會摘要 2026/06）")


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Theme 3: oil master-switch + rates red line
# ════════════════════════════════════════════════════════════════
s = content_slide("主題三：油價是「總開關」，利率是「紅線」")
prose(s, 0.3, 1.35, 7.55, 5.7, [
    "若說盈餘決定市場的高度，那麼**油價與利率，就決定了市場的下檔風險**。德銀直言油價是這輪行情的「總開關」：基準情境假設霍爾木茲海峽於 6 月底重啟、布倫特二季均價約 109 美元、四季回落至 **86 美元**；但若海峽封鎖拖延，油價恐逼近 **150 美元**，屆時全球增長風險將非線性放大，歐洲最先承壓。**6/16** 的美伊臨時協議使布倫特一度跌破 84 美元，但彭博提醒協議持久性偏低、海峽重啟仍脆弱。",
    "利率方面，高盛劃出明確紅線：**10 年期美債殖利率觸及 5%**，才是真正威脅股市估值的水準——目前尚未觸及，但債市已成為下半年最關鍵的變數。財政赤字接近 GDP 的 **6.6%**、長債供給增加、通膨回落放緩，使長端利率面臨上行壓力，殖利率曲線呈現「熊市平坦化」。真正的風險不在於 AI 故事結束，而在於**利率、槓桿與波動率開始相互放大**。",
    "兩者交匯出全場的勝負手——**油價能否回落到足以把 CPI 壓回 3% 以下**。若能，則花旗的降息劇本成立；若不能，則美銀所警告的「滯脹」路徑風險上升。",
], size=12.5, space_after=9)
card(s, 8.1, 1.5, 4.95, 1.5, "油價情境（布倫特，美元/桶）", GOLD_DK,
     ["基準 86　|　事故情境 150", "6/16 臨時協議後一度跌破 84"], body_size=11.5)
card(s, 8.1, 3.18, 4.95, 1.5, "利率紅線", NAVY,
     ["10 年期美債 5% ＝股市紅線", "財政赤字 ≈ GDP 6.6%，曲線熊平"], body_size=11.5)
card(s, 8.1, 4.86, 4.95, 1.5, "全場勝負手", UWRED,
     ["油價能否壓 CPI 回 3% 以下", "→ 花旗降息 vs 美銀滯脹"], body_size=11.5)
source(s, "資料來源：Deutsche Bank / Goldman Sachs / Citi / BofA / Bloomberg（PA 晨會摘要 2026/06）")


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Regional divergence + market temperature
# ════════════════════════════════════════════════════════════════
s = content_slide("區域分化與市場溫度")
prose(s, 0.3, 1.32, 12.7, 1.65, [
    "在共同的宏觀背景下，各區域表現明顯分化：**美國**最具韌性、是各大投行一致的超配市場，隱憂在估值與擁擠；**歐洲**最為脆弱，距技術性衰退僅一步之遙；**亞洲（除日）**為亮點，**韓國**受惠記憶體晶片超級週期、估值仍低，存在重估空間；**日本**被油價改寫、日圓走貶；**中國**資金面偏弱。與此同時，支撐行情的情緒與資金本身已成風險——擁擠度與樂觀程度雙雙逼近極端。",
], size=12, space_after=6)
table(s, 0.3, 3.15, 7.45, [
    ["區域", "觀點", "關鍵數據"],
    ["美國", "最具韌性，一致超配", "消費+設備投資；估值偏高"],
    ["亞洲(除日)", "亮點，韓國突圍", "MSCI 韓股 7.2x；IT EPS 翻倍"],
    ["日本", "被油價改寫", "日圓 161.8（24/7 來最低）"],
    ["歐洲", "最弱", "2026 GDP 0.5%；近技術性衰退"],
    ["中國", "資金偏弱", "連 11 週淨流出"],
], widths=[1.25, 2.4, 3.8], row_h=0.5, body_size=10.5, head_size=11)
# temperature stat chips (right) 2 x 3
stats = [
    ("39%", "AI 股佔標普權重", GOLD_DK),
    ("80%", "做多半導體＝史上最擁擠", UWRED),
    ("9.2", "美銀牛熊指標（賣出區）", UWRED),
    ("4.1%", "平均現金水位（防禦↑）", NAVY),
    ("$119bn", "美股單週淨流入（紀錄）", GOLD_DK),
    ("28%", "視 AI 泡沫為首要尾部風險", NAVY),
]
bx, by = 7.95, 3.15
cw, chh = 2.5, 1.12
for i, (num, lab, col) in enumerate(stats):
    x = bx + (i % 2) * (cw + 0.08)
    y = by + (i // 2) * (chh + 0.1)
    box(s, x, y, cw, chh, CARD)
    box(s, x, y, 0.06, chh, col)
    tf = tbox(s, x + 0.16, y + 0.08, cw - 0.24, 0.5).text_frame
    para(tf, md(num), 19, col, first=True, bold=True, space_after=0)
    tf2 = tbox(s, x + 0.16, y + 0.6, cw - 0.24, 0.45).text_frame
    para(tf2, md(lab), 9.5, INK, first=True, space_after=0, line=1.05)
source(s, "資料來源：Deutsche Bank / Morgan Stanley / BofA / Bloomberg（PA 晨會摘要 2026/06）")


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Conclusion + asset allocation
# ════════════════════════════════════════════════════════════════
s = content_slide("結論與資產配置建議")
prose(s, 0.3, 1.32, 12.7, 1.4, [
    "綜合過去四週各大投行觀點，我們維持**「建設性、但防禦微調」**的基調：基本面（盈餘）仍足以支撐市場，但在擁擠度創高、利率逼近紅線的環境下，應主動**降低集中與槓桿、提高現金緩衝**，並緊盯油價與利率兩大搖擺因子。具體配置上：",
], size=12.5, space_after=6)
cy, ch2 = 2.72, 3.85
card(s, 0.3, cy, 4.05, ch2, "▲  增持 OVERWEIGHT", OWGREEN, [
    "美股（品質/盈餘）：金融・工業・醫療，廣度外擴",
    "亞洲（除日）/ 韓國：晶片週期 + 低估值重估",
    "中期美債：殖利率年底回落，carry + 防禦",
    "新興市場低波動、高股息個股",
], body_size=11.5)
card(s, 4.64, cy, 4.05, ch2, "▼  減持 UNDERWEIGHT", UWRED, [
    "擁擠的大型科技 / 半導體集中部位（降槓桿）",
    "長天期債券：熊陡風險，10Y 逼近 5%",
    "歐洲股票 / 信用：近技術性衰退",
    "日圓 / 未對沖日股：結構性走弱",
], body_size=11.5)
card(s, 8.98, cy, 4.05, ch2, "●  中性 / 戰術 NEUTRAL", GOLD_DK, [
    "歐洲能源股：Hormuz 重啟 → 降至中性",
    "黃金：滯脹 / 地緣對沖，小幅配置",
    "現金：提高至 ~4%，保留乾火藥",
    "波動率資產：適度配置以防尾部",
], body_size=11.5)
tf = tbox(s, 0.3, 6.74, 12.7, 0.38).text_frame
para(tf, md("※ 上述為彙整各投行研究之整合視角，非個別投資建議；花旗 / 大摩對聯儲路徑（仍見降息）與多數機構（轉向加息）存在分歧，請依自身風險屬性判讀。"),
     9, SLATE, first=True, space_after=0)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Watch list / risk map
# ════════════════════════════════════════════════════════════════
s = content_slide("下一步：關鍵觀察點與風險地圖")
prose(s, 0.3, 1.32, 12.7, 1.0, [
    "未來數週，市場走向「1996 軟著陸」或「1999 過熱」劇本，取決於以下六個變數。任一變數惡化，都可能改變上述配置的攻守平衡——建議以此作為動態調整的觸發點：",
], size=12.5, space_after=4)
watch = [
    ("油價 → CPI", "能否壓回 3% 以下？花旗降息劇本 vs 美銀滯脹劇本的勝負手", GOLD_DK),
    ("FOMC 路徑", "10 月加息已定價；美銀看 3 碼、德銀 2 碼，花旗逆勢看降息", UWRED),
    ("10 年期美債", "逼近 5% ＝股市紅線；曲線熊陡、財政赤字與長端供給壓力", NAVY),
    ("AI 盈餘 / 資本開支", "2027 雲廠商 capex 看破 1 兆；盈餘可持續性是估值關鍵", GOLD_DK),
    ("擁擠度去化", "做多半導體 / Mag 7 若鬆動，恐引發連鎖去槓桿與波動放大", NAVY),
    ("政治 / 日圓", "11 月中期選舉「股債匯三殺」風險；日圓 carry 平倉尾部", UWRED),
]
gx, gy = 0.3, 2.5
gw, gh = 4.13, 1.78
for i, (t, d, col) in enumerate(watch):
    x = gx + (i % 3) * (gw + 0.07)
    y = gy + (i // 3) * (gh + 0.12)
    box(s, x, y, gw, gh, CARD)
    box(s, x, y, gw, 0.07, col)
    box(s, x, y, 0.07, gh, col)
    tf = tbox(s, x + 0.2, y + 0.14, gw - 0.34, 0.5).text_frame
    para(tf, md(t), 13.5, col, first=True, bold=True, space_after=0)
    tf2 = tbox(s, x + 0.2, y + 0.66, gw - 0.36, gh - 0.78).text_frame
    para(tf2, md(d), 11, INK, first=True, space_after=0, line=1.18)
box(s, 0.3, 6.62, 12.73, 0.5, NAVY)
tf = tbox(s, 0.45, 6.66, 12.4, 0.42, MSO_ANCHOR.MIDDLE).text_frame
para(tf, md("底線：基本面仍穩，但市場容錯率低 —— 增持盈餘確定性、控制集中與槓桿，緊盯油價與利率兩大搖擺因子。"),
     12.5, WHITE, first=True, bold=True, space_after=0)


force_all_fonts()

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
